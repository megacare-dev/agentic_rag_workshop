#!/usr/bin/env python3
"""
Generate PowerPoint: Agentic RAG Workshop — Pre-Session
Theme: Mahidol University (Navy Blue + Gold)

Output: pre_session_slides.pptx
Upload to Google Slides → File → Import slides → select file

Usage: python3 create_slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Mahidol University Theme Colors ─────────────────────────────────────────
NAVY       = RGBColor(0x1A, 0x23, 0x70)   # Deep Mahidol Blue
DARK_NAVY  = RGBColor(0x0F, 0x14, 0x52)   # Darker background
GOLD       = RGBColor(0xC4, 0x99, 0x1C)   # Mahidol Gold
LIGHT_GOLD = RGBColor(0xEA, 0xC7, 0x50)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xDD)
CODE_BG    = RGBColor(0x1E, 0x1E, 0x2E)   # Dark code bg
CODE_GREEN = RGBColor(0x4E, 0xC9, 0x63)   # Green code text
CODE_BLUE  = RGBColor(0x7A, 0xA2, 0xF7)   # Blue keywords
CODE_ORANGE= RGBColor(0xFA, 0xB3, 0x87)   # Orange strings
FAINT_GOLD = RGBColor(0x8B, 0x7A, 0x40)   # Subtle gold for accents


def set_slide_bg(slide, color=NAVY):
    """Set slide background to solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height):
    """Add a textbox and return the text frame."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def add_run(paragraph, text, size=16, bold=False, color=WHITE,
            font_name="Sarabun", italic=False):
    """Add a styled run to a paragraph."""
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return run


def new_para(tf, text="", size=16, bold=False, color=WHITE,
             font_name="Sarabun", align=PP_ALIGN.LEFT, space_after=Pt(6),
             space_before=Pt(0)):
    """Add a new paragraph with a run."""
    p = tf.add_paragraph()
    p.alignment = align
    p.space_after = space_after
    p.space_before = space_before
    if text:
        add_run(p, text, size, bold, color, font_name)
    return p


def add_gold_line(slide, top_inches, left=1.0, width=8.0):
    """Add a horizontal gold accent line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top_inches),
        Inches(width), Pt(3),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()
    return shape


def add_code_block(slide, code_text, left, top, width, height):
    """Add a styled code block with dark background."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)

    # Clear default paragraph
    tf.paragraphs[0].text = ""

    lines = code_text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(1)
        p.space_before = Pt(1)

        # Simple syntax highlighting
        run = p.add_run()
        run.text = line
        run.font.size = Pt(10)
        run.font.name = "Roboto Mono"

        # Color based on content
        stripped = line.strip()
        if stripped.startswith("#"):
            run.font.color.rgb = RGBColor(0x6A, 0x99, 0x55)  # Green comments
        elif any(stripped.startswith(kw) for kw in
                 ["import ", "from ", "for ", "try:", "except ", "if ", "else:", "break", "raise", "def "]):
            run.font.color.rgb = CODE_BLUE
        elif "'" in stripped or '"' in stripped:
            run.font.color.rgb = CODE_ORANGE
        else:
            run.font.color.rgb = CODE_GREEN

    return shape


def add_table(slide, rows_data, left, top, width, height):
    """Add a styled table."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])

    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(left), Inches(top),
        Inches(width), Inches(height),
    )
    table = table_shape.table

    # Set column widths proportionally
    total_w = Inches(width)
    col_widths = [total_w // n_cols] * n_cols
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for r, row in enumerate(rows_data):
        for c, cell_text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = ""

            # Set cell background
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_NAVY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY if r % 2 == 0 else RGBColor(0x15, 0x1B, 0x60)

            # Set text
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = cell_text
            run.font.size = Pt(12)
            run.font.name = "Sarabun"
            run.font.bold = (r == 0)
            run.font.color.rgb = GOLD if r == 0 else WHITE

    return table_shape


def add_callout(slide, text, top=4.6):
    """Add a gold callout bar at the bottom."""
    # Background bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(top),
        Inches(9.0), Inches(0.45),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x22, 0x2B, 0x80)
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_run(p, text, size=13, bold=True, color=LIGHT_GOLD)

    return shape


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════════

def slide_01_cover(prs):
    """Cover / Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, DARK_NAVY)

    add_gold_line(slide, 1.7)
    add_gold_line(slide, 3.8)

    # Title
    tf = add_textbox(slide, 1.0, 0.8, 8.0, 1.5)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    add_run(p, "Agentic RAG", size=48, bold=True, color=WHITE)
    p2 = new_para(tf, "From Zero to Hero", size=36, bold=True, color=GOLD,
                  space_before=Pt(8))

    # Subtitle
    tf2 = add_textbox(slide, 1.0, 4.0, 8.0, 1.0)
    new_para(tf2, "Workshop 4 ชั่วโมง", size=18, color=LIGHT_GRAY,
             space_after=Pt(4))
    # Remove the default first empty paragraph
    tf2.paragraphs[0].text = ""
    p = tf2.paragraphs[0]
    add_run(p, "Workshop 4 ชั่วโมง", size=18, color=LIGHT_GRAY)
    new_para(tf2, "คณะ ICT มหาวิทยาลัยมหิดล", size=18, color=LIGHT_GOLD)


def slide_02_learning(prs):
    """สิ่งที่จะได้เรียนรู้"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_gold_line(slide, 0.85, left=0.7, width=3.5)

    tf = add_textbox(slide, 0.7, 0.3, 8.5, 0.6)
    p = tf.paragraphs[0]
    add_run(p, "สิ่งที่จะได้เรียนรู้", size=32, bold=True, color=GOLD)

    items = [
        ("✂️", "เตรียมข้อมูลสำหรับ RAG (Chunk → Embed → VectorDB)"),
        ("🤖", "สร้าง AI Agent ด้วย Google ADK"),
        ("⭐", "สร้าง Agentic RAG — Agent ที่ค้นหา + คิด + ตอบเอง"),
        ("📊", "วัดคุณภาพด้วย LLM-as-Judge"),
    ]
    tf2 = add_textbox(slide, 1.0, 1.3, 7.5, 3.5)
    tf2.paragraphs[0].text = ""
    for emoji, text in items:
        p = tf2.paragraphs[0] if tf2.paragraphs[0].text == "" else tf2.add_paragraph()
        p.space_after = Pt(18)
        add_run(p, f"{emoji}  ", size=22, color=GOLD)
        add_run(p, text, size=20, color=WHITE)


def slide_03_timeline(prs):
    """Timeline"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_gold_line(slide, 0.85, left=0.7, width=3.0)

    tf = add_textbox(slide, 0.7, 0.3, 8.5, 0.6)
    p = tf.paragraphs[0]
    add_run(p, "Timeline วัน Workshop", size=32, bold=True, color=GOLD)

    table_data = [
        ["Part", "เวลา", "เนื้อหา"],
        ["📢 Part 1", "1 ชม. 20 นาที", "Data Pipeline: Chunk → Embed → Qdrant"],
        ["☕ พัก", "10 นาที", ""],
        ["📢 Part 2", "1 ชม. 30 นาที", "Agent → Tool → RAG Agent → Agentic RAG"],
        ["🧪 Part 3", "1 ชม.", "แบบฝึกหัด (10 คะแนน) + Q&A"],
    ]
    add_table(slide, table_data, 0.7, 1.3, 8.5, 3.0)


def slide_04_register(prs):
    """ลงทะเบียน"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_gold_line(slide, 0.8, left=0.7, width=5.5)

    tf = add_textbox(slide, 0.7, 0.2, 8.5, 0.6)
    p = tf.paragraphs[0]
    add_run(p, "🌐 ลงทะเบียน — myhero.megawiz.co.th", size=28, bold=True, color=GOLD)

    # Left: Steps
    tf2 = add_textbox(slide, 0.7, 1.0, 4.5, 2.8)
    p = tf2.paragraphs[0]
    add_run(p, "ขั้นตอนลงทะเบียน", size=16, bold=True, color=LIGHT_GOLD)

    steps = [
        "1. เข้า myhero.megawiz.co.th/student-portal",
        "2. กรอก อีเมล เพื่อรับ OTP",
        "3. ใส่รหัส OTP → เข้าสู่ระบบ",
        "4. กรอก ชื่อ-นามสกุล และ รหัสนักศึกษา",
    ]
    for s in steps:
        new_para(tf2, s, size=14, color=WHITE, space_after=Pt(6))

    # Right: Benefits
    tf3 = add_textbox(slide, 5.3, 1.0, 4.2, 2.8)
    p = tf3.paragraphs[0]
    add_run(p, "สิ่งที่จะได้จากระบบ", size=16, bold=True, color=LIGHT_GOLD)

    benefits = [
        "📄 เอกสารประกอบการเรียน (Notebook)",
        "📝 แบบฝึกหัด + ส่งงาน",
        "🏆 Certificate หลังจบ workshop",
    ]
    for b in benefits:
        new_para(tf3, b, size=14, color=WHITE, space_after=Pt(8))

    add_callout(slide, "🎯 ลงทะเบียนก่อนมาเรียน เพื่อให้วัน workshop ไม่ต้องเสียเวลา")


def slide_05_prep(prs):
    """สิ่งที่ต้องเตรียม"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_gold_line(slide, 0.85, left=0.7, width=5.0)

    tf = add_textbox(slide, 0.7, 0.3, 8.5, 0.6)
    p = tf.paragraphs[0]
    add_run(p, "⚠️ สิ่งที่ต้องเตรียมก่อนมาเรียน", size=32, bold=True, color=GOLD)

    items = [
        ("✅", "ลงทะเบียน ที่ myhero.megawiz.co.th/student-portal"),
        ("🔑", "สร้าง Gemini API Key"),
        ("💻", "เปิด Google Colab ได้"),
        ("🧪", "ทดสอบ API Key ล่วงหน้า"),
    ]
    tf2 = add_textbox(slide, 1.0, 1.4, 8.0, 3.0)
    tf2.paragraphs[0].text = ""
    for i, (emoji, text) in enumerate(items):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.space_after = Pt(18)
        add_run(p, f"{i+1}. {emoji}  ", size=22, bold=True, color=LIGHT_GOLD)
        add_run(p, text, size=20, color=WHITE)


def slide_06_apikey(prs):
    """วิธีสร้าง Gemini API Key"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_gold_line(slide, 0.8, left=0.7, width=4.5)

    tf = add_textbox(slide, 0.7, 0.2, 8.5, 0.6)
    p = tf.paragraphs[0]
    add_run(p, "🔑 วิธีสร้าง Gemini API Key", size=28, bold=True, color=GOLD)

    steps = [
        'เข้า aistudio.google.com',
        'Login ด้วย Google Account',
        'คลิก "Get API Key" (มุมบนซ้าย)',
        'คลิก "Create API Key"',
        'เลือก Project → Create API key in new project',
        'Copy API Key เก็บไว้ (ห้ามแชร์!)',
    ]
    tf2 = add_textbox(slide, 0.7, 1.0, 8.5, 3.2)
    tf2.paragraphs[0].text = ""
    for i, step in enumerate(steps):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.space_after = Pt(10)
        add_run(p, f"{i+1}. ", size=16, bold=True, color=LIGHT_GOLD)
        add_run(p, step, size=16, color=WHITE)

    add_callout(slide, "⚠️ Free tier: 15 RPM (requests/minute) — เพียงพอสำหรับ workshop")


def slide_07_secrets(prs):
    """เก็บ API Key ใน Colab Secrets"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_gold_line(slide, 0.8, left=0.7, width=5.5)

    tf = add_textbox(slide, 0.7, 0.2, 8.5, 0.6)
    p = tf.paragraphs[0]
    add_run(p, "🔒 วิธีเก็บ API Key ใน Colab Secrets", size=28, bold=True, color=GOLD)

    steps = [
        'เปิด Google Colab',
        'คลิก 🔑 ไอคอนกุญแจ (แถบซ้าย) → "Secrets"',
        'คลิก "Add new secret"',
        'Name: GOOGLE_API_KEY',
        'Value: วาง API Key ที่ copy มา',
        'เปิด toggle "Notebook access" ✅',
    ]
    tf2 = add_textbox(slide, 0.7, 1.0, 8.5, 3.2)
    tf2.paragraphs[0].text = ""
    for i, step in enumerate(steps):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.space_after = Pt(10)
        add_run(p, f"{i+1}. ", size=16, bold=True, color=LIGHT_GOLD)
        add_run(p, step, size=16, color=WHITE)

    add_callout(slide, "🎯 ทำแค่ครั้งเดียว — ใช้ได้ทุก notebook")


def slide_08_ratelimit(prs):
    """Error 429"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_gold_line(slide, 0.8, left=0.7, width=4.0)

    tf = add_textbox(slide, 0.7, 0.2, 8.5, 0.6)
    p = tf.paragraphs[0]
    add_run(p, "🚨 Error 429 — Rate Limit คืออะไร?", size=28, bold=True, color=GOLD)

    code = "google.api_core.exceptions.ResourceExhausted: 429\nYou exceeded your current quota"
    add_code_block(slide, code, 0.7, 1.0, 8.5, 0.7)

    tf2 = add_textbox(slide, 0.7, 1.9, 8.5, 0.5)
    p = tf2.paragraphs[0]
    add_run(p, "แปลว่า: ", size=16, bold=True, color=GOLD)
    add_run(p, "ส่ง request เร็ว/เยอะเกินไป", size=16, color=WHITE)

    # Free tier limits table
    tf3 = add_textbox(slide, 0.7, 2.5, 8.5, 0.4)
    p = tf3.paragraphs[0]
    add_run(p, "Free Tier Limits (Gemini 2.5 Flash):", size=14, bold=True, color=LIGHT_GOLD)

    table_data = [
        ["Limit", "ค่า"],
        ["RPM (Requests/min)", "15"],
        ["TPM (Tokens/min)", "1,000,000"],
        ["RPD (Requests/day)", "1,500"],
    ]
    add_table(slide, table_data, 0.7, 3.0, 8.5, 2.0)


def slide_09_fix429(prs):
    """วิธีแก้ Error 429"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_gold_line(slide, 0.8, left=0.7, width=3.0)

    tf = add_textbox(slide, 0.7, 0.2, 8.5, 0.6)
    p = tf.paragraphs[0]
    add_run(p, "💡 วิธีแก้ Error 429", size=28, bold=True, color=GOLD)

    # Quick fix section
    tf2 = add_textbox(slide, 0.7, 0.9, 4.0, 1.5)
    p = tf2.paragraphs[0]
    add_run(p, "แก้ทันที:", size=16, bold=True, color=LIGHT_GOLD)
    p2 = new_para(tf2, "1. รอ 1-2 นาที แล้ว run ใหม่", size=13, color=WHITE, space_after=Pt(6))
    p3 = new_para(tf2, "2. อย่า run cell ซ้ำเร็วๆ — รอ output ก่อน", size=13, color=WHITE)

    # Permanent fix label
    tf3 = add_textbox(slide, 5.0, 0.9, 4.5, 0.4)
    p = tf3.paragraphs[0]
    add_run(p, "แก้ถาวร (แนะนำ):", size=16, bold=True, color=LIGHT_GOLD)

    # Code block
    code = (
        "import time\n"
        "for attempt in range(3):\n"
        "    try:\n"
        "        response = client.models.generate_content(...)\n"
        "        break\n"
        "    except Exception as e:\n"
        "        if '429' in str(e):\n"
        "            print(f'⏳ Rate limit, รอ {30*(attempt+1)} วินาที...')\n"
        "            time.sleep(30 * (attempt + 1))\n"
        "        else:\n"
        "            raise"
    )
    add_code_block(slide, code, 0.5, 2.5, 9.0, 2.5)


def slide_10_test(prs):
    """ทดสอบก่อนมาเรียน"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_gold_line(slide, 0.8, left=0.7, width=3.5)

    tf = add_textbox(slide, 0.7, 0.2, 8.5, 0.6)
    p = tf.paragraphs[0]
    add_run(p, "🧪 ทดสอบก่อนมาเรียน", size=28, bold=True, color=GOLD)

    tf2 = add_textbox(slide, 0.7, 0.75, 8.5, 0.3)
    p = tf2.paragraphs[0]
    add_run(p, "เปิด Colab → สร้าง notebook ใหม่ → Run:", size=13, color=LIGHT_GRAY)

    code = (
        "# 1. ตั้ง API Key\n"
        "from google.colab import userdata\n"
        "import os\n"
        "os.environ['GOOGLE_API_KEY'] = userdata.get('GOOGLE_API_KEY')\n"
        "\n"
        "# 2. ทดสอบ\n"
        "from google import genai\n"
        "client = genai.Client(api_key=os.environ['GOOGLE_API_KEY'])\n"
        "resp = client.models.generate_content(\n"
        "    model='gemini-3.1-pro-preview',\n"
        "    contents='สวัสดี ตอบสั้นๆ 1 ประโยค'\n"
        ")\n"
        "print(f'✅ สำเร็จ: {resp.text}')"
    )
    add_code_block(slide, code, 0.7, 1.1, 8.5, 2.8)

    add_callout(slide, "ถ้าเห็น ✅ = พร้อมเรียน!", top=4.2)


def slide_11_checklist(prs):
    """Checklist"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_gold_line(slide, 0.85, left=0.7, width=3.5)

    tf = add_textbox(slide, 0.7, 0.3, 8.5, 0.6)
    p = tf.paragraphs[0]
    add_run(p, "📋 Checklist ก่อนมาเรียน", size=32, bold=True, color=GOLD)

    items = [
        "ลงทะเบียนที่ myhero.megawiz.co.th/student-portal แล้ว",
        "มี Google Account",
        "สร้าง Gemini API Key แล้ว",
        "เก็บ Key ใน Colab Secrets แล้ว",
        "Run ทดสอบผ่าน → เห็น ✅",
    ]
    tf2 = add_textbox(slide, 1.0, 1.3, 8.0, 3.0)
    tf2.paragraphs[0].text = ""
    for i, item in enumerate(items):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.space_after = Pt(14)
        add_run(p, "☐  ", size=22, bold=True, color=LIGHT_GOLD)
        add_run(p, item, size=18, color=WHITE)

    add_callout(slide, "🎯 ทำ checklist ครบ → วัน workshop จะราบรื่นมาก!")


def slide_12_cert(prs):
    """Certificate"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_gold_line(slide, 0.85, left=0.7, width=2.5)

    tf = add_textbox(slide, 0.7, 0.3, 8.5, 0.6)
    p = tf.paragraphs[0]
    add_run(p, "🏆 Certificate", size=32, bold=True, color=GOLD)

    tf2 = add_textbox(slide, 0.7, 1.2, 8.5, 0.4)
    p = tf2.paragraphs[0]
    add_run(p, "หลังจบ workshop และส่งแบบฝึกหัดครบ:", size=16, color=LIGHT_GRAY)

    items = [
        ("🤖", "ระบบตรวจให้คะแนนอัตโนมัติ (AI Grading)"),
        ("📜", "รับ Certificate ผ่านระบบ myhero.megawiz.co.th"),
        ("📥", "ดาวน์โหลดได้ทันที (PDF)"),
    ]
    tf3 = add_textbox(slide, 1.0, 1.8, 8.0, 2.5)
    tf3.paragraphs[0].text = ""
    for i, (emoji, text) in enumerate(items):
        p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
        p.space_after = Pt(18)
        add_run(p, f"{emoji}  ", size=24, color=GOLD)
        add_run(p, text, size=20, color=WHITE)


def slide_13_help(prs):
    """มีปัญหา?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_gold_line(slide, 0.75, left=0.7, width=2.0)

    tf = add_textbox(slide, 0.7, 0.2, 8.5, 0.5)
    p = tf.paragraphs[0]
    add_run(p, "❓ มีปัญหา?", size=32, bold=True, color=GOLD)

    table_data = [
        ["ปัญหา", "วิธีแก้"],
        ["เข้า myhero ไม่ได้", "ตรวจอีเมล spam / ลองใหม่"],
        ["ไม่ได้รับ OTP", "เช็ค spam folder / รอ 1 นาที"],
        ["API key not valid", "สร้าง key ใหม่ที่ aistudio.google.com"],
        ["Error 429", "รอ 1-2 นาที แล้วลองใหม่"],
        ["Colab ช้า", "Runtime → Change runtime type → T4 GPU"],
    ]
    add_table(slide, table_data, 0.7, 0.9, 8.5, 3.2)

    add_callout(slide, "ถามอาจารย์ / TA ได้เลย 🙋", top=4.4)


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    # Build all 13 slides
    slide_01_cover(prs)
    slide_02_learning(prs)
    slide_03_timeline(prs)
    slide_04_register(prs)
    slide_05_prep(prs)
    slide_06_apikey(prs)
    slide_07_secrets(prs)
    slide_08_ratelimit(prs)
    slide_09_fix429(prs)
    slide_10_test(prs)
    slide_11_checklist(prs)
    slide_12_cert(prs)
    slide_13_help(prs)

    output_path = "pre_session_slides.pptx"
    prs.save(output_path)
    print(f"✅ สร้างไฟล์สำเร็จ: {output_path}")
    print(f"📤 Upload ไปยัง Google Slides:")
    print(f"   1. เข้า slides.google.com")
    print(f"   2. คลิก + (Blank) → File → Import slides")
    print(f"   3. หรือ upload ไฟล์ .pptx ไปที่ Google Drive แล้วเปิดด้วย Google Slides")


if __name__ == "__main__":
    main()
